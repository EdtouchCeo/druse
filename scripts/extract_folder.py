# -*- coding: utf-8 -*-
"""폴더 내 PDF/PPTX/HWPX 자료를 하나의 검색용 raw md로 통합 추출.

사용법:
  python scripts/extract_folder.py <원본폴더> <출력md경로> <제목>

예:
  python scripts/extract_folder.py "input/teacher/법령지침" "data/법령지침_raw.md" "법령 · 지침 검색 데이터"
  python scripts/extract_folder.py "input/teacher/교육과정" "data/교육과정_raw.md" "교육과정 검색 데이터"

필요 패키지: pymupdf(fitz), python-pptx, lxml
"""
import sys, os, zipfile, re
from lxml import etree

def extract_pdf(path):
    import fitz
    doc = fitz.open(path)
    out = []
    for i in range(len(doc)):
        t = doc[i].get_text("text").strip()
        if t:
            out.append(f"[p.{i+1}]\n{t}")
    doc.close()
    return "\n\n".join(out)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for si, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if txt:
                    parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))
        if parts:
            out.append(f"[슬라이드 {si}]\n" + "\n".join(parts))
    return "\n\n".join(out)

def extract_hwpx(path):
    """HWPX = ZIP. Contents/section*.xml 안의 텍스트를 문단 단위로 추출."""
    out = []
    with zipfile.ZipFile(path) as z:
        names = [n for n in z.namelist() if re.match(r"Contents/section\d+\.xml$", n)]
        names.sort(key=lambda n: int(re.search(r"section(\d+)", n).group(1)))
        for n in names:
            data = z.read(n)
            root = etree.fromstring(data)
            paras = []
            # 문단(localname=p) 단위로, 그 안의 텍스트(localname=t) 결합
            for p in root.iter():
                if etree.QName(p).localname != "p":
                    continue
                texts = []
                for el in p.iter():
                    ln = etree.QName(el).localname
                    if ln == "t" and el.text:
                        texts.append(el.text)
                line = "".join(texts).strip()
                if line:
                    paras.append(line)
            if paras:
                out.append("\n".join(paras))
    return "\n\n".join(out)

# 일부 PDF는 어절 사이 공백을 폭 없는 문자(ZWNJ 등)로 넣는다 — 그대로 두면
# "심의위원회의<ZWNJ>위원은"처럼 붙어 보이고 키워드 검색도 걸리지 않는다. 공백으로 되돌린다.
ZERO_WIDTH_RE = re.compile("[​‌‍﻿]+")


def normalize(text):
    return ZERO_WIDTH_RE.sub(" ", text or "")


DISPATCH = {".pdf": extract_pdf, ".pptx": extract_pptx, ".hwpx": extract_hwpx}

# 여기서 뽑지 않는 파일 — 글자가 이미지·도형이라 추출하면 0자가 나오고,
# 대신 쪽을 판독한 문면을 data/sources/*.json에 두어 build_law_extras.py가 덧붙인다.
# 이 목록에 없으면 빈 문서 자리표시자 청크가 검색 인덱스에 섞인다.
SKIP = {
    "(붙임2) 선생님을 위한 교육활동 보호 가이드 리플렛.pdf":
        "data/sources/dge_teacher_protection_guide.json",
}


def main(folder, out_path, title):
    import unicodedata

    def skipped(name):
        n = unicodedata.normalize("NFC", name)   # OneDrive NFD 파일명 대응
        return any(unicodedata.normalize("NFC", k) == n for k in SKIP)

    files = sorted(f for f in os.listdir(folder)
                   if os.path.splitext(f)[1].lower() in DISPATCH and not skipped(f))
    chunks = [f"# {title}\n",
              f"> 검색용 raw 데이터 — 원본 폴더: `{folder}`",
              f"> 총 {len(files)}개 문서\n",
              "## 수록 문서 목록\n"]
    for i, f in enumerate(files, 1):
        chunks.append(f"{i}. {f}")
    chunks.append("\n---\n")
    report = []
    for f in files:
        ext = os.path.splitext(f)[1].lower()
        fp = os.path.join(folder, f)
        try:
            text = normalize(DISPATCH[ext](fp)).strip()
        except Exception as e:
            text = ""
            report.append((f, "ERROR: " + str(e)))
        if not text:
            report.append((f, "빈 텍스트(스캔본/추출불가 가능)"))
            text = "_(텍스트를 추출하지 못했습니다 — 스캔 이미지 PDF이거나 형식 미지원)_"
        else:
            report.append((f, f"{len(text):,}자"))
        chunks.append(f"\n\n# 📄 {f}\n\n{text}")
    result = "\n".join(chunks)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as w:
        w.write(result)
    print(f"== {title} ==")
    print(f"출력: {out_path}  ({len(result):,}자)")
    for name, info in report:
        print(f"  - {name}: {info}")

if __name__ == "__main__":
    if len(sys.argv) != 4:
        print(__doc__); sys.exit(1)
    main(sys.argv[1], sys.argv[2], sys.argv[3])
